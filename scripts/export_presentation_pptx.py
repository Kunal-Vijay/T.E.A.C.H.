#!/usr/bin/env python3
"""Export T.E.A.C.H. hackathon HTML deck to PowerPoint (.pptx)."""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PRESENTATION = ROOT / "presentation"
OUT_PATH = PRESENTATION / "TEACH-Hackathon.pptx"
DEMO_DATA = PRESENTATION / "demo-data.json"
MEME_GIF = PRESENTATION / "giphy.gif"
NOVA_IMG = ROOT / "teach-frontend/public/image-from-rawpixel-id-12165579-png.png"

# Light theme palette (matches styles.css)
BG = RGBColor(244, 247, 251)
TEXT = RGBColor(15, 23, 42)
MUTED = RGBColor(100, 116, 139)
TEAL = RGBColor(13, 148, 136)
TEAL_DARK = RGBColor(15, 118, 110)
AMBER = RGBColor(217, 119, 6)
WHITE = RGBColor(255, 255, 255)


def load_demo_data() -> dict:
    if DEMO_DATA.exists():
        return json.loads(DEMO_DATA.read_text(encoding="utf-8"))
    return {
        "stats": {"topics": 0, "sessions": 0, "teach_sessions": 0, "doubt_sessions": 0, "viva_sessions": 0},
        "topics": [],
        "featured_topic": {"title": "Newton Laws", "toc": ["First Law", "Second Law"]},
    }


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.LEFT,
    font_name: str = "Inter",
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_eyebrow(slide, text: str, top=Inches(0.45)) -> None:
    add_textbox(slide, Inches(0.6), top, Inches(10), Inches(0.35), text.upper(), size=11, bold=True, color=TEAL_DARK)


def add_title(slide, text: str, top=Inches(0.85), size: int = 32) -> None:
    add_textbox(slide, Inches(0.6), top, Inches(11.5), Inches(1.4), text, size=size, bold=True, color=TEXT)


def add_body(slide, text: str, top, height=Inches(1.2), size: int = 14, color: RGBColor = MUTED) -> None:
    add_textbox(slide, Inches(0.6), top, Inches(11.5), height, text, size=size, color=color)


def add_bullets(slide, items: list[str], left, top, width, height, size: int = 13) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(6)
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = TEXT
        run.font.name = "Inter"


def slide_1_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Logo block
    logo = slide.shapes.add_shape(1, Inches(0.6), Inches(0.55), Inches(0.85), Inches(0.85))  # rectangle
    logo.fill.solid()
    logo.fill.fore_color.rgb = TEAL
    logo.line.fill.background()
    add_textbox(slide, Inches(0.72), Inches(0.68), Inches(0.6), Inches(0.5), "T", size=28, bold=True, color=WHITE)

    add_textbox(slide, Inches(0.6), Inches(1.55), Inches(11), Inches(0.9), "T.E.A.C.H.", size=44, bold=True)
    add_textbox(
        slide,
        Inches(0.6),
        Inches(2.45),
        Inches(11),
        Inches(0.8),
        "A concept gap filler — between every student and their teacher.",
        size=22,
        color=TEAL_DARK,
    )
    add_body(
        slide,
        "Teacherless Education through Autonomous Cognitive Heuristics",
        Inches(3.25),
        height=Inches(0.5),
        size=13,
    )
    add_body(slide, "Infinity Learn · Hackathon 2026", Inches(3.75), height=Inches(0.4), size=12)

    team = (
        "Nikhil Kumar Singh · Anishwar Sharma · Divyanshu Kumar · "
        "Kunal Vijay · Abhishek Kumar Singh · Vikrant Pandey"
    )
    add_body(slide, team, Inches(4.35), height=Inches(0.8), size=11)


def slide_2_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "The Problem")
    add_title(
        slide,
        "Our teachers are exceptional.\nBut they can't be available every time, everywhere.",
        size=26,
    )
    add_body(
        slide,
        "Batch size, class time, and student hesitation — gaps still grow after the bell rings.",
        Inches(2.35),
        height=Inches(0.5),
    )

    journey_lines = [
        "IN CLASS — \"I understood everything.\" 😎",
        "↓",
        "LATER — \"Wait… what was Newton's Second Law?\" 💀",
        "↓",
        "STUDENT — \"If I ask the teacher again… what will they think?\"",
        "↓",
        "NOVA — \"Let's start from the basics.\"",
    ]
    add_bullets(slide, journey_lines, Inches(0.6), Inches(2.95), Inches(5.8), Inches(3.2), size=12)

    if MEME_GIF.exists():
        slide.shapes.add_picture(str(MEME_GIF), Inches(2.0), Inches(4.55), width=Inches(2.8))

    cards = [
        "👥 Batch-level teaching — Exceptional teachers still can't clear every doubt in limited class time.",
        "🤐 Student hesitation — \"If I go again and again… what will they think?\"",
        "⏱ Concept gaps linger — Rarely time to re-teach one concept for the whole batch.",
    ]
    add_bullets(slide, cards, Inches(6.7), Inches(2.95), Inches(5.5), Inches(3.5), size=11)


def slide_3_solution(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "Our Solution")
    add_title(slide, "A concept gap filler that bridges student and teacher.", size=26)
    add_body(
        slide,
        "Our teachers are exceptional — Nova doesn't replace them. When a student needs help "
        "after class, at home, or late at night, Nova acts as a bridge so no concept gap is left behind.",
        Inches(2.15),
        height=Inches(1.0),
    )
    add_textbox(
        slide,
        Inches(0.6),
        Inches(3.2),
        Inches(6),
        Inches(0.5),
        "Student  ↔  Nova · Gap Filler  ↔  Teacher",
        size=16,
        bold=True,
        color=TEAL_DARK,
    )
    add_body(
        slide,
        "Missed a formula? Lagging on a concept? Ask Nova — anytime, without judgment.",
        Inches(3.65),
        height=Inches(0.5),
    )

    dialogue = [
        "Student: \"Can you explain this?\"",
        "Nova: \"Of course.\"",
        "Student: \"…again?\"",
        "Nova: \"Of course.\"",
    ]
    add_bullets(slide, dialogue, Inches(7.2), Inches(0.9), Inches(5.2), Inches(1.5), size=11)

    if NOVA_IMG.exists():
        slide.shapes.add_picture(str(NOVA_IMG), Inches(8.0), Inches(2.4), width=Inches(2.2))

    caps = ["Explains", "Adapts", "Teaches", "Revises"]
    add_bullets(slide, caps, Inches(7.2), Inches(5.0), Inches(5.2), Inches(1.2), size=12)


def slide_4_how_it_helps(prs: Presentation, demo: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "How It Helps")
    add_title(slide, "Not just watching. Not just listening.\nLearning with feedback.", size=26)

    modes = [
        "TEACH ME — Fill concept gaps with step-by-step whiteboard teaching.",
        "ASK A DOUBT — \"What was that formula from today's class?\"",
        "VOICE VIVA 👀 — Light spoken viva; Nova asks so you know where you stand.",
    ]
    add_bullets(slide, modes, Inches(0.6), Inches(2.2), Inches(11.5), Inches(1.4), size=12)

    s = demo["stats"]
    stats = (
        f"Live catalog · teach.db — {s['topics']} topics · {s['sessions']} sessions · "
        f"{s['teach_sessions']} teach · {s['doubt_sessions']} doubt · {s['viva_sessions']} viva"
    )
    add_body(slide, stats, Inches(3.55), height=Inches(0.4), size=11, color=TEAL_DARK)

    topic_lines = [f"{t['subject']}: {t['title']}" for t in demo.get("topics", [])]
    add_bullets(slide, topic_lines, Inches(0.6), Inches(4.0), Inches(11.5), Inches(1.2), size=10)

    impact = [
        "• Learner progress tracking",
        "• Higher engagement — active, not passive",
        "• Summarisation potential (roadmap)",
        "• Voice narration · whiteboard · adaptive pace",
        "• \"I just have ONE small doubt.\" → full teaching session. 🚀",
    ]
    add_bullets(slide, impact, Inches(0.6), Inches(5.2), Inches(11.5), Inches(1.5), size=11)


def slide_5_why_it_works(prs: Presentation, demo: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "Why It Works")
    add_title(slide, "AI can answer. Nova teaches.", size=28)

    left = [
        "WATCH & READ ALONE",
        "Student: \"I'll rewatch the lecture.\"",
        "Student: \"Still confused. I'll ask tomorrow.\" 💀",
    ]
    ft = demo.get("featured_topic", {})
    toc = ", ".join(ft.get("toc", ["First Law", "Second Law"]))
    right = [
        "NOVA · GAP FILLER",
        "Student: \"I didn't get Newton's Second Law from class.\"",
        "Nova: \"Imagine you're sitting in a moving bus…\"",
        f"Nova: \"Build the concept — {toc} — and check with a viva.\"",
    ]
    add_bullets(slide, left, Inches(0.6), Inches(1.75), Inches(5.5), Inches(2.2), size=11)
    add_bullets(slide, right, Inches(6.4), Inches(1.75), Inches(5.8), Inches(2.4), size=11)

    tech = "React + Vite · FastAPI · AWS Bedrock · Claude · Amazon Nova Sonic · Google Cloud TTS · Structured Prompts"
    add_body(slide, tech, Inches(4.2), height=Inches(0.5), size=10)

    s = demo["stats"]
    add_body(
        slide,
        f"Built on teach.db — {s['topics']} published topics, {s['sessions']} real sessions. Demo: {ft.get('title', 'Newton Laws')}.",
        Inches(4.65),
        height=Inches(0.5),
        size=10,
        color=TEAL_DARK,
    )


def slide_6_closing(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "T.E.A.C.H.")
    add_textbox(
        slide,
        Inches(0.6),
        Inches(1.5),
        Inches(11.5),
        Inches(2.0),
        "Our teachers are exceptional.\nWe believe every student deserves access to one — anytime.",
        size=30,
        bold=True,
        color=TEXT,
    )
    add_body(
        slide,
        "Nova fills the gaps when teachers can't be everywhere.\n"
        "Teachers stay at the centre. Learning never stops because the classroom did.",
        Inches(3.75),
        height=Inches(1.2),
        size=16,
    )


def main() -> None:
    demo = load_demo_data()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_1_title(prs)
    slide_2_problem(prs)
    slide_3_solution(prs)
    slide_4_how_it_helps(prs, demo)
    slide_5_why_it_works(prs, demo)
    slide_6_closing(prs)

    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
